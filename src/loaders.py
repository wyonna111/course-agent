"""多格式文档加载：PDF / PPT / TXT，支持图片多模态描述"""

import io
from pathlib import Path

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document

from src.config import MIN_TEXT_LEN


def _load_pdf(file_path: Path) -> list[Document]:
    from src.vision import describe_image, is_vision_enabled
    vision_on = is_vision_enabled()

    docs = PyPDFLoader(str(file_path)).load()
    result: list[Document] = []
    for d in docs:
        text = (d.page_content or "").strip()
        page = d.metadata.get("page", len(result))
        if len(text) < MIN_TEXT_LEN:
            text = text or "（本页未提取到可检索文字，可能为图片或空白页）"
        d.page_content = text
        d.metadata["doc_type"] = "pdf"
        d.metadata["source_name"] = file_path.name
        d.metadata["display_page"] = int(page) + 1
        d.metadata["has_image"] = False
        result.append(d)

    if len(result) >= 2:
        # 用 fitz 补充图片描述（如果视觉开启）
        if vision_on:
            try:
                import fitz
                pdf = fitz.open(str(file_path))
                for i, page_obj in enumerate(pdf):
                    image_descs = []
                    for img in page_obj.get_images(full=True):
                        xref = img[0]
                        try:
                            img_bytes = pdf.extract_image(xref)["image"]
                            if len(img_bytes) > 5000:
                                desc = describe_image(img_bytes)
                                if desc:
                                    image_descs.append(desc)
                        except Exception:
                            pass
                    if image_descs and i < len(result):
                        result[i].page_content += "\n[图片内容：" + "；".join(image_descs) + "]"
                        result[i].metadata["has_image"] = True
                pdf.close()
            except ImportError:
                pass
        return result

    try:
        import fitz

        pdf = fitz.open(str(file_path))
        docs = []
        for i, page_obj in enumerate(pdf):
            text = (page_obj.get_text() or "").strip()
            if len(text) < MIN_TEXT_LEN:
                text = text or "（本页未提取到可检索文字，可能为图片或空白页）"

            image_descs = []
            if vision_on:
                for img in page_obj.get_images(full=True):
                    xref = img[0]
                    try:
                        img_bytes = pdf.extract_image(xref)["image"]
                        if len(img_bytes) > 5000:
                            desc = describe_image(img_bytes)
                            if desc:
                                image_descs.append(desc)
                    except Exception:
                        pass
            if image_descs:
                text += "\n[图片内容：" + "；".join(image_descs) + "]"

            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "page": i,
                        "display_page": i + 1,
                        "source": str(file_path),
                        "source_name": file_path.name,
                        "doc_type": "pdf",
                        "has_image": len(image_descs) > 0,
                    },
                )
            )
        pdf.close()
        if docs:
            return docs
    except ImportError:
        pass

    if not docs:
        raise ValueError(
            f"{file_path.name} 几乎读不到文字（可能是扫描版 PDF）。"
            "请换可复制文字的 PDF 或导出为 .txt"
        )
    for d in docs:
        d.metadata.setdefault("doc_type", "pdf")
        d.metadata.setdefault("source_name", file_path.name)
    return docs


def _extract_pptx_images(slide) -> list[bytes]:
    """提取一张幻灯片内的所有图片字节，过滤过小的图标。"""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
            try:
                img_bytes = shape.image.blob
                if len(img_bytes) > 5000:  # 过滤小图标（<5KB）
                    images.append(img_bytes)
            except Exception:
                pass
    return images


def _load_pptx(file_path: Path) -> list[Document]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("请安装 python-pptx：pip install python-pptx") from e

    from src.vision import describe_image, is_vision_enabled
    vision_on = is_vision_enabled()

    prs = Presentation(str(file_path))
    docs: list[Document] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        text = "\n".join(parts).strip()
        if len(text) < MIN_TEXT_LEN:
            text = text or "（本幻灯片未提取到可检索文字，可能为图片页）"

        # 图片理解：提取图片 → 生成描述 → 追加到正文
        image_descs: list[str] = []
        if vision_on:
            for img_bytes in _extract_pptx_images(slide):
                desc = describe_image(img_bytes)
                if desc:
                    image_descs.append(desc)

        if image_descs:
            text += "\n[图片内容：" + "；".join(image_descs) + "]"

        docs.append(
            Document(
                page_content=text,
                metadata={
                    "page": i,
                    "slide": i + 1,
                    "display_page": i + 1,
                    "source": str(file_path),
                    "source_name": file_path.name,
                    "doc_type": "pptx",
                    "has_image": len(image_descs) > 0,
                },
            )
        )
    if not docs:
        raise ValueError(f"{file_path.name} 中未读到有效文字（可能多为图片幻灯片）")
    return docs


def _load_text(file_path: Path) -> list[Document]:
    docs = TextLoader(str(file_path), encoding="utf-8").load()
    for d in docs:
        d.metadata["doc_type"] = "text"
        d.metadata["source_name"] = file_path.name
        d.metadata["page"] = 0
    return docs


def load_documents(file_path: Path) -> list[Document]:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(file_path)
    if suffix in {".pptx", ".ppt"}:
        return _load_pptx(file_path)
    if suffix in {".txt", ".md"}:
        return _load_text(file_path)
    raise ValueError(f"不支持的格式：{suffix}，目前支持 PDF / PPTX / TXT / MD")


def location_label(meta: dict) -> str:
    """生成可溯源位置标签，供引用展示。"""
    name = meta.get("source_name") or Path(meta.get("source", "未知文件")).name
    doc_type = meta.get("doc_type", "file")
    if doc_type == "pptx":
        slide = meta.get("slide") or int(meta.get("page", 0)) + 1
        return f"[本地: {name} 幻灯片 {slide}]"
    page = meta.get("page")
    if page is not None:
        return f"[本地: {name} 第 {int(page) + 1} 页]"
    return f"[本地: {name}]"
